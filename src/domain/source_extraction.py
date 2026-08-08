"""Extract searchable text from uploaded knowledge sources."""
from __future__ import annotations

import csv
import io
import json
import os
import re
import socket
import struct
import subprocess
import zipfile
from functools import lru_cache
from pathlib import Path
from typing import Any
import structlog
from src.core.config import settings

logger = structlog.get_logger()

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".xlsx", ".xlsm", ".pptx", ".txt", ".md", ".csv",
    ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp",
}


class SourceExtractionError(ValueError):
    pass


def _scan_with_clamd(data: bytes) -> None:
    """Scan bytes through ClamAV's TCP INSTREAM protocol without temp files."""
    host = settings.MALWARE_SCANNER_HOST
    if not host:
        raise SourceExtractionError("Malware scanning is unavailable")
    try:
        with socket.create_connection((host, settings.MALWARE_SCANNER_PORT), timeout=10) as client:
            client.settimeout(30)
            client.sendall(b"zINSTREAM\0")
            for offset in range(0, len(data), 1024 * 1024):
                chunk = data[offset:offset + 1024 * 1024]
                client.sendall(struct.pack("!I", len(chunk)) + chunk)
            client.sendall(struct.pack("!I", 0))
            response = client.recv(4096).decode("utf-8", errors="replace")
    except OSError as exc:
        raise SourceExtractionError("Malware scanning is unavailable") from exc
    if "OK" not in response or "FOUND" in response:
        raise SourceExtractionError("The uploaded file failed malware scanning")


def _page(number: int, text: str) -> dict[str, Any]:
    return {"page_number": number, "text": _clean(text)}


def _clean(text: str) -> str:
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _token_coverage(original: str, candidate: str) -> float:
    original_tokens = set(re.findall(r"[A-Za-zÀ-ỹ0-9][A-Za-zÀ-ỹ0-9_-]{3,}", original.lower()))
    if not original_tokens:
        return 1.0
    candidate_tokens = set(re.findall(r"[A-Za-zÀ-ỹ0-9][A-Za-zÀ-ỹ0-9_-]{3,}", candidate.lower()))
    return len(original_tokens & candidate_tokens) / len(original_tokens)


def _validate_archive(data: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            members = archive.infolist()
            if len(members) > settings.MAX_SOURCE_ARCHIVE_FILES:
                raise SourceExtractionError("The document archive contains too many files")
            total_size = sum(max(0, item.file_size) for item in members)
            if total_size > settings.MAX_SOURCE_UNCOMPRESSED_BYTES:
                raise SourceExtractionError("The document archive expands beyond the allowed size")
            for item in members:
                if item.compress_size and item.file_size / item.compress_size > 10_000:
                    raise SourceExtractionError("The document archive has an unsafe compression ratio")
    except zipfile.BadZipFile as exc:
        raise SourceExtractionError("The uploaded document is not a valid archive") from exc


def _validate_source_bytes(filename: str, data: bytes) -> None:
    extension = Path(filename).suffix.lower()
    if extension in {".docx", ".xlsx", ".xlsm", ".pptx"}:
        if not data.startswith(b"PK"):
            raise SourceExtractionError("The uploaded Office document has an invalid file signature")
        _validate_archive(data)
    elif extension == ".pdf" and not data.lstrip().startswith(b"%PDF"):
        raise SourceExtractionError("The uploaded PDF has an invalid file signature")
    elif extension in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp"}:
        try:
            from PIL import Image
            Image.MAX_IMAGE_PIXELS = settings.MAX_SOURCE_IMAGE_PIXELS
            with Image.open(io.BytesIO(data)) as image:
                image.verify()
        except Exception as exc:
            raise SourceExtractionError("The uploaded image is invalid or unsafe") from exc

    if settings.MALWARE_SCAN_ENABLED:
        if settings.MALWARE_SCANNER_HOST:
            _scan_with_clamd(data)
        else:
            try:
                result = subprocess.run(
                    [settings.MALWARE_SCANNER_COMMAND, "--stream", "--no-summary"],
                    input=data,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=30,
                    check=False,
                )
            except (FileNotFoundError, subprocess.TimeoutExpired) as exc:
                raise SourceExtractionError("Malware scanning is unavailable") from exc
            if result.returncode != 0:
                raise SourceExtractionError("The uploaded file failed malware scanning")


@lru_cache(maxsize=1)
def _markitdown() -> Any:
    try:
        from markitdown import MarkItDown
    except Exception:
        return None
    try:
        return MarkItDown(enable_plugins=False)
    except TypeError:
        # MarkItDown 0.0.x has no plugin constructor argument. Its built-in
        # converters are still sufficient here because PaddleOCR remains the
        # scanned-page fallback.
        return MarkItDown()


def _convert_with_markitdown(filename: str, data: bytes) -> str:
    if not settings.MARKITDOWN_ENABLED:
        return ""
    converter = _markitdown()
    if converter is None:
        return ""
    try:
        result = converter.convert_stream(
            io.BytesIO(data),
            file_extension=Path(filename).suffix.lower(),
        )
        value = getattr(result, "markdown", None) or getattr(result, "text_content", None) or ""
        return _clean(str(value))
    except Exception:
        # MarkItDown is an enhancement layer. A single unsupported or malformed
        # file must still be handled by the existing format-specific extractor.
        return ""


@lru_cache(maxsize=1)
def _paddle_ocr() -> Any:
    # PaddlePaddle's Windows oneDNN executor currently fails on OCR model
    # attributes; disable that optional accelerator and use the CPU path.
    os.environ.setdefault("FLAGS_use_mkldnn", "0")
    os.environ.setdefault("FLAGS_use_onednn", "0")
    try:
        import paddle
        paddle.set_flags({"FLAGS_use_mkldnn": False, "FLAGS_use_onednn": False})
        from paddleocr import PaddleOCR
    except Exception as exc:
        raise SourceExtractionError(
            "PaddleOCR is not available. Install paddlepaddle and paddleocr to process scanned files."
        ) from exc
    try:
        return PaddleOCR(
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            lang=settings.PADDLEOCR_LANG,
            enable_mkldnn=False,
        )
    except TypeError:
        return PaddleOCR(use_angle_cls=True, lang=settings.PADDLEOCR_LANG, show_log=False, enable_mkldnn=False)


def _ocr_image(image: Any) -> str:
    engine = _paddle_ocr()
    # PaddleOCR 3.x accepts NumPy arrays (not PIL objects) and returns a
    # Result whose JSON payload nests recognition fields under ``res``.
    try:
        import numpy as np
        image_input = np.asarray(image) if not isinstance(image, (str, np.ndarray)) else image
    except Exception:
        image_input = image

    def collect_texts(payload: Any) -> list[str]:
        if callable(payload):
            payload = payload()
        if isinstance(payload, str):
            try:
                payload = json.loads(payload)
            except json.JSONDecodeError:
                return []
        if not isinstance(payload, dict):
            return []
        result = payload.get("res", payload)
        if not isinstance(result, dict):
            return []
        return [str(value) for value in result.get("rec_texts", []) if value]

    lines: list[str] = []
    if hasattr(engine, "predict"):
        for item in engine.predict(image_input):
            payload = item.json if hasattr(item, "json") else item
            lines.extend(collect_texts(payload))
        if lines:
            return _clean("\n".join(lines))

    try:
        legacy_result = engine.ocr(image_input, cls=True)
    except TypeError:
        # PaddleOCR 3.x keeps ``ocr`` as a compatibility alias but removed
        # the ``cls`` keyword; PaddleOCR 2.x still needs it.
        legacy_result = engine.ocr(image_input)
    for page in legacy_result or []:
        for line in page or []:
            try:
                lines.append(str(line[1][0]))
            except (IndexError, TypeError):
                continue
    return _clean("\n".join(lines))


def _extract_pdf_pages(data: bytes) -> list[dict[str, Any]]:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = [_page(index, page.extract_text() or "") for index, page in enumerate(reader.pages, start=1)]
        if any(item["text"] for item in pages):
            # Mixed PDFs are common: retain embedded text and OCR only image
            # pages instead of silently dropping scanned appendices.
            if any(not item["text"] for item in pages):
                try:
                    import fitz
                    from PIL import Image
                    document = fitz.open(stream=data, filetype="pdf")
                    for index, item in enumerate(pages):
                        if item["text"]:
                            continue
                        page = document[index]
                        if not page.get_images(full=True):
                            continue
                        pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
                        item["text"] = _ocr_image(Image.open(io.BytesIO(pixmap.tobytes("png"))))
                except Exception:
                    # Embedded-text PDFs remain usable when OCR is not
                    # installed; blank pages are preserved rather than
                    # converting the entire document to a failing request.
                    pass
            return pages
    except Exception:
        pass

    try:
        import fitz
        from PIL import Image
        document = fitz.open(stream=data, filetype="pdf")
        pages: list[dict[str, Any]] = []
        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            pages.append(_page(index, _ocr_image(Image.open(io.BytesIO(pixmap.tobytes("png"))))))
        return pages
    except SourceExtractionError:
        raise
    except Exception as exc:
        raise SourceExtractionError(f"Could not extract text from PDF: {exc}") from exc


def _extract_docx(data: bytes) -> str:
    from docx import Document
    document = Document(io.BytesIO(data))
    parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    return _clean("\n".join(parts))


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in workbook.worksheets:
        parts.append(f"## {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                parts.append(" | ".join(values))
    return _clean("\n".join(parts))


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation
    presentation = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"## Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return _clean("\n".join(parts))


def extract_source_pages(filename: str, data: bytes) -> list[dict[str, Any]]:
    extension = Path(filename).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise SourceExtractionError(
            f"Unsupported file type '{extension or 'unknown'}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    if not data:
        raise SourceExtractionError("The uploaded file is empty.")
    _validate_source_bytes(filename, data)
    if extension == ".pdf":
        pages = _extract_pdf_pages(data)
    elif extension == ".docx":
        pages = [_page(1, _extract_docx(data))]
    elif extension in {".xlsx", ".xlsm"}:
        pages = [_page(1, _extract_xlsx(data))]
    elif extension == ".pptx":
        pages = [_page(1, _extract_pptx(data))]
    elif extension == ".csv":
        rows = csv.reader(io.StringIO(data.decode("utf-8-sig", errors="replace")))
        pages = [_page(1, "\n".join(" | ".join(cell.strip() for cell in row) for row in rows))]
    elif extension in {".txt", ".md"}:
        pages = [_page(1, data.decode("utf-8-sig", errors="replace"))]
    else:
        from PIL import Image
        import numpy as np
        Image.MAX_IMAGE_PIXELS = settings.MAX_SOURCE_IMAGE_PIXELS
        image = Image.open(io.BytesIO(data)).convert("RGB")
        pages = [_page(1, _ocr_image(np.asarray(image)))]
    pages = [item for item in pages if item["text"]]
    if len(pages) > settings.MAX_SOURCE_PAGES:
        raise SourceExtractionError(f"Documents are limited to {settings.MAX_SOURCE_PAGES} pages")
    if sum(len(str(item["text"])) for item in pages) > settings.MAX_SOURCE_TEXT_CHARS:
        raise SourceExtractionError("The extracted document text is too large")
    if not pages:
        raise SourceExtractionError("No readable text was found in the uploaded file.")
    return pages


def extract_source_markdown(
    filename: str,
    data: bytes,
    pages: list[dict[str, Any]] | None = None,
) -> str:
    """Convert a source to Markdown without weakening page-aware extraction.

    MarkItDown supplies document structure (headings, lists, tables and links).
    ``pages`` remains the authoritative page-indexed OCR/text representation
    used for citations and original-source review. If MarkItDown returns an
    incomplete result, the page extraction is used instead.
    """
    fallback = _clean("\n\n".join(str(item["text"]) for item in (pages or extract_source_pages(filename, data))))
    converted = _convert_with_markitdown(filename, data)
    if not converted:
        logger.info("Source Markdown conversion used page extractor", filename=filename, reason="markitdown_unavailable_or_failed")
        return fallback
    coverage = _token_coverage(fallback, converted)
    if len(converted) < max(80, int(len(fallback) * 0.30)) or coverage < 0.80:
        logger.warning(
            "MarkItDown output rejected; using page extractor",
            filename=filename,
            fallback_characters=len(fallback),
            markdown_characters=len(converted),
            token_coverage=round(coverage, 3),
        )
        return fallback
    logger.info(
        "Source converted to Markdown with MarkItDown",
        filename=filename,
        page_extractor_characters=len(fallback),
        markdown_characters=len(converted),
        token_coverage=round(coverage, 3),
    )
    return converted


def extract_source(filename: str, data: bytes) -> str:
    """Return the legacy flattened representation used by article bodies."""
    return _clean("\n\n".join(str(item["text"]) for item in extract_source_pages(filename, data)))
